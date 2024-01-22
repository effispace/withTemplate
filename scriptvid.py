from pptx import Presentation
from pptx.util import Inches
import subprocess
import fitz  # PyMuPDF
import os
from moviepy.editor import *
from pptx_to_pdf import convert_pptx_to_pdf
import re

#importing API key etc...
import openai
from openai import OpenAI
from dotenv import load_dotenv

load_dotenv()
openai_api = os.getenv('OPENAI_API')

if openai_api is None:
    raise ValueError("The OPEN_AI environment variable is not set.")
#don't change this part of the code. Should remain the same

#Set Global Variables
powerpoint_file_path = "my_presentation.pptx"   # Path to the PowerPoint file
pdf_file_path = "my_presentation.pdf"  # Path to the PDF file
images_folder = "slide_images"  # Path to the folder where the images will be stored
video_file_path = "my_presentation.mp4"  # Path to the output video file
frame_rate = 30  # Adjust the frame rate as needed
duration_per_image = 3  # Time each image is shown in seconds
transition_duration = 0.5  # Duration of the transition between images in seconds

template_file = 'template.pptx'#load the powerpoint template
prs = Presentation(template_file)

#generate chatgpt slides content LLM content generation (prompt argument is the topic of the presentation)
def create_content(prompt):
    print("Converting to educational content....")
    client = openai.OpenAI(api_key=openai_api)
    story = client.chat.completions.create(
        messages = [
        {"role": "system", "content": "Create content for a 20-slide presentation on the input topic. For each slide, provide a title and up to six bullet points that effectively cover different aspects of the topic. Ensure that the content is informative, accurate, and well-structured to facilitate a comprehensive understanding of the subject. Lable the title of each slide as Title: followed by the name of the title of the slide"},
        {"role": "user", "content": prompt}],
        model="gpt-3.5-turbo",
    )
    # write response to the file
    nextpass = story.choices[0].message.content  
    file_name = "chatgptresponse.txt"
    with open(file_name, 'w') as file:
        # Write the string to the file
        file.write(nextpass)
def add_title(titleName):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = titleName

def addContent(titleName,contentBox):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = titleName
    content = slide.placeholders[1]
    content.text = contentBox
 
def imageContent(imgpath, titleName, contentBox,slideType):
    slide_layout = prs.slide_layouts[slideType]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = titleName
    content = slide.placeholders[1]
    content.text = contentBox
    image_placeholder = slide.placeholders[2]
    image_path = imgpath
    left = image_placeholder.left
    top = image_placeholder.top
    width = image_placeholder.width
    height = image_placeholder.height
    slide.shapes.add_picture(image_path, left, top, width, height)

#create presentation based on chatgpt text (stored in a text file)
def create_presentation_from_text(file_path):
    with open(file_path, 'r') as file:
        content = file.read()
    slides = re.split(r'Title:', content)[1:]

    for slide in slides:
        title_match = re.match(r'(.+)', slide)
        bullet_points_match = re.findall(r'- (.+)', slide)
        if title_match:
            title = title_match.group(1).strip()
            bullet_points = '\n'.join(bullet_points_match) 
            addContent(title, bullet_points)  
    add_title("Thanks for listening <3")
    prs.save(powerpoint_file_path)

def pdf_to_images(pdf_path, output_folder):
    pdf_document = fitz.open(pdf_path)
    
    for page_num in range(pdf_document.page_count):
        page = pdf_document[page_num]
        image = page.get_pixmap()
        
        image_path = f"{output_folder}/page_{page_num + 1}.png"
        image.save(image_path)

    pdf_document.close()

def create_clip(image_path, clip_duration):
    return ImageClip(image_path, duration=clip_duration)

def get_images(images_folder):
    # Get all the images path in the folder
    images = [img for img in os.listdir(images_folder) if img.endswith(".png")]

    def get_page_number(image_name):
        # Extract the numerical part of the filename
        return int(''.join(filter(str.isdigit, image_name)))

    images.sort(key=get_page_number)

    images = [os.path.join(images_folder, img) for img in images]

    return images

def create_slideshow(images_folder, output_video_path, frame_rate, duration_per_image, transition_duration):

    images = get_images(images_folder)

    image_clips = [create_clip(image, duration_per_image) for image in images]

    def fadeout(clip, duration):
        #return the portion of the clip when it is fading out
        fade = clip.subclip(clip.duration - duration, clip.duration)
        return CompositeVideoClip([fade.crossfadeout(duration)])
    
    def fadein(clip, duration):
        #return the portion of the clip when it is fading in
        fade = clip.subclip(0, duration)
        return CompositeVideoClip([fade.crossfadein(duration)])

    final_clip = []

    for i in range(0, len(image_clips)):
        clip = image_clips[i]
        fadein_clip = fadein(clip, transition_duration)
        fadeout_clip = fadeout(clip, transition_duration)
        final_clip.append(fadein_clip)
        final_clip.append(clip.set_duration(clip.duration - transition_duration*2))
        final_clip.append(fadeout_clip)

    concatenate_videoclips(final_clip).write_videofile(output_video_path, codec='libx264', fps=frame_rate)

if __name__ == "__main__":
    topic = "galaxies"
    create_content(topic) #generate chatgpt script for presentation - parameter is topic
    #change prompt as you feel is right in order to get a script that will flow well 
   
    file_path = "chatgptresponse.txt"
    add_title(topic)
    create_presentation_from_text(file_path)# create slides from script file 

    convert_pptx_to_pdf(powerpoint_file_path, pdf_file_path)

    pdf_to_images(pdf_file_path, images_folder)

    create_slideshow(images_folder, video_file_path, frame_rate, duration_per_image, transition_duration)